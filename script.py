from pptx import Presentation

def overwrite_pptx_with_thanks(file_path):
    presentation = Presentation()

    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    # Agrego el texto "Gracias" en el centro de la diapositiva
    textbox = slide.shapes.add_textbox(
        left=presentation.slide_width // 4,
        top=presentation.slide_height // 3,
        width=presentation.slide_width // 2,
        height=presentation.slide_height // 3
    )
    text_frame = textbox.text_frame
    text_frame.text = "Gracias"

    presentation.save(file_path)

# Ruta del archivo
file_path = "trial.pptx"
overwrite_pptx_with_thanks(file_path)

print(f"El archivo {file_path} ha sido sobrescrito con el mensaje 'Gracias'.")
